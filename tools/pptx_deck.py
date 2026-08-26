# -*- coding: utf-8 -*-
"""非機能要件 ヒアリング資料 生成
既存 doc/powerpoint/要件定義.pptx のデザイントークンを踏襲する。
"""
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- design tokens (既存 pptx から抽出) ----
NAVY  = RGBColor(0x1F, 0x38, 0x64)
GREEN = RGBColor(0x1E, 0x7B, 0x4D)
INK   = RGBColor(0x1D, 0x21, 0x25)
GRAY  = RGBColor(0x5A, 0x64, 0x72)
LINE  = RGBColor(0xE4, 0xE9, 0xF0)
ZEBRA = RGBColor(0xF7, 0xF9, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xB4, 0x53, 0x09)
FONT  = "Yu Gothic UI"

L = 0.8888888888888888          # 左マージン
W = 11.5555                     # コンテンツ幅
TOP_TITLE, TOP_BADGE, TOP_SUB, TOP_LEAD = 0.33, 0.44, 0.79, 1.11
TABLE_TOP = 1.75
BOTTOM = 7.12


def _tb(slide, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def text(slide, x, y, w, h, s, size=11, bold=False, color=INK, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, space=0.0):
    tf = _tb(slide, x, y, w, h)
    tf.vertical_anchor = anchor
    for i, ln in enumerate(str(s).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space:
            p.space_before = Pt(space)
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = color
        p.line_spacing = 1.25
    return tf


def rect(slide, x, y, w, h, fill):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(slide, title, badge=None, sub=None, lead=None):
    rect(slide, 0, 0, 0.12, 1.03, NAVY)
    text(slide, L, TOP_TITLE, 9.6, 0.44, title, size=24, bold=True, color=NAVY)
    if badge:
        text(slide, 10.5, TOP_BADGE, 1.94, 0.31, badge, size=14, bold=True,
             color=GRAY, align=PP_ALIGN.RIGHT)
    if sub:
        text(slide, L, TOP_SUB, W, 0.28, sub, size=12, color=GRAY)
    if lead:
        cpl = int((W - 0.2) / (13 * 0.0139))          # 13pt での 1 行文字数
        n = max(1, math.ceil(len(lead) / cpl))
        h = max(0.30, n * 13 * 1.30 / 72.0 + 0.06)
        rect(slide, L, TOP_LEAD, 0.04, h, GREEN)
        text(slide, L + 0.195, TOP_LEAD - 0.02, W - 0.2, h, lead, size=13, bold=True, color=INK)
        return TOP_LEAD + h + 0.24
    return TOP_LEAD


# ---- 行テーブル ----
# cols: list of dict(x, w, label, size, bold, color)
def rowtable(slide, top, cols, rows, pad=0.15):
    rect(slide, L, top, W, 0.36, NAVY)
    for c in cols:
        text(slide, c["x"], top + 0.085, c["w"], 0.22, c["label"], size=11, bold=True, color=WHITE)
    y = top + 0.36
    heights = []
    for row in rows:
        h = 0.0
        for c, cell in zip(cols, row):
            size = c.get("size", 11)
            cpl = max(4, int(c["w"] / (size * 0.0139)))     # 全角換算の1行文字数
            lines = 0
            for ln in str(cell).split("\n"):
                lines += max(1, math.ceil(len(ln) / cpl))
            h = max(h, lines * size * 1.30 / 72.0)
        heights.append(h + pad * 2)
    total = sum(heights)
    avail = BOTTOM - y
    if total > avail:                                        # 収まらなければ均等圧縮
        k = avail / total
        heights = [h * k for h in heights]
    for i, (row, h) in enumerate(zip(rows, heights)):
        if i % 2 == 1:
            rect(slide, L, y, W, h, ZEBRA)
        rect(slide, L, y + h, W, 0.01, LINE)
        for c, cell in zip(cols, row):
            text(slide, c["x"], y + pad - 0.03, c["w"], h - pad, cell,
                 size=c.get("size", 11), bold=c.get("bold", False),
                 color=c.get("color", INK))
        y += h
    return y


# 標準4列（設問 / Q番号 / 選択肢 / 未確定だと何が決まらないか）
COLS = [
    dict(x=1.00, w=2.28, label="ご確認いただきたいこと", size=12, bold=True, color=INK),
    dict(x=3.40, w=0.62, label="設問",  size=11, color=GRAY),
    dict(x=4.12, w=3.52, label="選択肢（【推】＝当方の推奨・現行の暫定値）", size=10.5, color=GRAY),
    dict(x=7.78, w=4.60, label="未確定だと何が決まらないか", size=10.5, color=INK),
]


# ============================================================
# 概念図まわりの描画ヘルパ
# ============================================================
RED    = RGBColor(0xB0, 0x3A, 0x2E)
F_NAVY = RGBColor(0xF0, 0xF4, 0xFA)
F_GRN  = RGBColor(0xF2, 0xF8, 0xF4)
F_AMB  = RGBColor(0xFF, 0xF8, 0xEE)
F_RED  = RGBColor(0xFD, 0xF3, 0xF2)
TINT   = {NAVY: F_NAVY, GREEN: F_GRN, AMBER: F_AMB, RED: F_RED, GRAY: ZEBRA}


def card(slide, x, y, w, h, accent=NAVY, edge="top"):
    """淡い塗り＋アクセント罫のカード。"""
    rect(slide, x, y, w, h, TINT.get(accent, ZEBRA))
    if edge == "top":
        rect(slide, x, y, w, 0.05, accent)
    elif edge == "left":
        rect(slide, x, y, 0.05, h, accent)
    return x, y, w, h


def arrow(slide, x, y, w, h, color=GRAY, direction="right"):
    from pptx.enum.shapes import MSO_SHAPE
    shape = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def step(slide, x, y, w, h, no, title, body, accent=NAVY):
    """フロー図の 1 ステップ。"""
    card(slide, x, y, w, h, accent)
    text(slide, x + 0.16, y + 0.20, w - 0.32, 0.22, no, size=10, bold=True, color=accent)
    text(slide, x + 0.16, y + 0.46, w - 0.32, 0.50, title, size=12.5, bold=True, color=INK)
    text(slide, x + 0.16, y + h - 0.16 - _lines(body, w - 0.32, 10) * 10 * 1.30 / 72.0,
         w - 0.32, _lines(body, w - 0.32, 10) * 10 * 1.30 / 72.0, body, size=10, color=GRAY)


def _lines(s, w, size):
    cpl = max(3, int(w / (size * 0.0139)))
    return sum(max(1, math.ceil(len(ln) / cpl)) for ln in str(s).split("\n"))


def band(slide, x, y, w, h, accent, label, mid, right):
    """横帯 1 本（層の説明）。"""
    card(slide, x, y, w, h, accent, edge="left")
    text(slide, x + 0.26, y + 0.16, 2.85, h - 0.3, label, size=13, bold=True, color=NAVY)
    text(slide, x + 3.30, y + 0.16, w - 3.30 - 2.55, h - 0.3, mid, size=11, color=INK)
    text(slide, x + w - 2.40, y + 0.16, 2.20, h - 0.3, right, size=11, bold=True, color=accent)


def note(slide, x, y, w, s, accent=AMBER, size=11):
    h = _lines(s, w - 0.5, size) * size * 1.30 / 72.0 + 0.30
    rect(slide, x, y, w, h, TINT.get(accent, ZEBRA))
    rect(slide, x, y, 0.05, h, accent)
    text(slide, x + 0.26, y + 0.15, w - 0.5, h - 0.3, s, size=size, bold=True, color=INK)
    return y + h


def divider(prs, kicker, title, sub):
    s = blank(prs)
    rect(s, 0, 0, 0.28, 7.5, NAVY)
    text(s, 1.2, 2.75, 10.6, 0.32, kicker, size=14, bold=True, color=GREEN)
    text(s, 1.2, 3.20, 10.6, 0.8, title, size=34, bold=True, color=NAVY)
    rect(s, 1.2, 4.30, 2.2, 0.05, GREEN)
    text(s, 1.2, 4.62, 10.6, 0.8, sub, size=14, color=INK)
    return s
